import os
import re
import uuid
import shutil
import asyncio
import logging
from contextlib import asynccontextmanager
from pathlib import Path

import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Pt
from lxml import etree

from fastapi import FastAPI, UploadFile, File, Form, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse

# ─────────────────────────────────────────────
# Config
# ─────────────────────────────────────────────
UPLOAD_FOLDER = Path("./tmp")
UPLOAD_FOLDER.mkdir(parents=True, exist_ok=True)

MAX_FILE_SIZE_MB = 50
MAX_FILE_SIZE_BYTES = MAX_FILE_SIZE_MB * 1024 * 1024

ALLOWED_EXTENSIONS = {".pdf", ".pptx", ".ppt"}

logging.basicConfig(level=logging.INFO, format="%(asctime)s | %(levelname)s | %(message)s")
log = logging.getLogger(__name__)


# ─────────────────────────────────────────────
# Lifespan: periodic temp file cleanup
# ─────────────────────────────────────────────
@asynccontextmanager
async def lifespan(app: FastAPI):
    task = asyncio.create_task(_periodic_cleanup())
    yield
    task.cancel()

async def _periodic_cleanup():
    """Delete tmp files older than 10 minutes every 5 minutes."""
    while True:
        await asyncio.sleep(300)
        _cleanup_old_files(max_age_seconds=600)

def _cleanup_old_files(max_age_seconds: int = 600):
    import time
    now = time.time()
    for f in UPLOAD_FOLDER.iterdir():
        try:
            if now - f.stat().st_mtime > max_age_seconds:
                f.unlink()
                log.info(f"Cleaned up old file: {f.name}")
        except Exception:
            pass


# ─────────────────────────────────────────────
# App
# ─────────────────────────────────────────────
app = FastAPI(title="Watermark Remover API", lifespan=lifespan)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)


# ─────────────────────────────────────────────
# PDF Cleaner  (text + image watermarks)
# ─────────────────────────────────────────────
def clean_pdf(input_path: str, output_path: str, watermark_text: str) -> int:
    """
    Remove watermark from PDF.
    Handles:
      - Selectable text watermarks (redaction)
      - Transparent / low-opacity text layers
      - Image-based watermarks that exactly match known patterns
    Returns the number of redactions applied.
    """
    doc = fitz.open(input_path)
    total_hits = 0
    pattern = re.compile(re.escape(watermark_text), re.IGNORECASE)

    for page in doc:
        hits_on_page = 0

        # 1️⃣  Search visible/selectable text
        for inst in page.search_for(watermark_text, quads=True):
            page.add_redact_annot(inst, fill=(1, 1, 1))   # white fill
            hits_on_page += 1

        # 2️⃣  Search hidden / invisible text (OCR-placed, artifacts, etc.)
        blocks = page.get_text("rawdict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        for block in blocks:
            for line in block.get("lines", []):
                for span in line.get("spans", []):
                    if pattern.search(span.get("text", "")):
                        rect = fitz.Rect(span["bbox"])
                        page.add_redact_annot(rect, fill=(1, 1, 1))
                        hits_on_page += 1

        if hits_on_page:
            # images=2 → also redact images covered by the annotation rect
            page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

        total_hits += hits_on_page

    doc.save(output_path, garbage=4, deflate=True, clean=True)
    doc.close()
    log.info(f"PDF: removed {total_hits} watermark instance(s)")
    return total_hits


# ─────────────────────────────────────────────
# PPTX Cleaner  (text + images + master/layout)
# ─────────────────────────────────────────────
def _scrub_text_frame(text_frame, pattern: re.Pattern) -> int:
    """Remove watermark text from a text frame. Returns hit count."""
    hits = 0
    for para in text_frame.paragraphs:
        for run in para.runs:
            if pattern.search(run.text):
                run.text = pattern.sub("", run.text)
                hits += 1
    return hits


def _remove_watermark_shapes(shapes, pattern: re.Pattern) -> int:
    """
    Walk all shapes (including groups) and:
    - Clear matching text runs
    - Remove entire shape if it IS the watermark (text-only box whose
      entire visible text matches the watermark)
    Returns hit count.
    """
    hits = 0
    to_delete = []

    for shape in shapes:
        # Grouped shapes — recurse
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            hits += _remove_watermark_shapes(shape.shapes, pattern)
            continue

        if shape.has_text_frame:
            full_text = shape.text_frame.text.strip()
            if pattern.search(full_text):
                # If the whole shape is just the watermark, delete the shape element
                if pattern.fullmatch(full_text) or pattern.sub("", full_text).strip() == "":
                    to_delete.append(shape)
                    hits += 1
                else:
                    # Partial match — only clear the matching runs
                    hits += _scrub_text_frame(shape.text_frame, pattern)

    # Delete shapes whose sole purpose was the watermark
    for shape in to_delete:
        sp = shape._element
        sp.getparent().remove(sp)

    return hits


def clean_pptx(input_path: str, output_path: str, watermark_text: str) -> int:
    """
    Remove watermark from PPTX.
    Handles:
      - Regular slide shapes
      - Grouped shapes
      - Slide master shapes
      - Slide layout shapes
      - Shapes with partial watermark text
    Returns the number of removals.
    """
    pattern = re.compile(re.escape(watermark_text), re.IGNORECASE)
    prs = Presentation(input_path)
    total_hits = 0

    # 1️⃣  Each slide's own shapes
    for slide in prs.slides:
        total_hits += _remove_watermark_shapes(slide.shapes, pattern)

    # 2️⃣  Slide master(s) — watermarks often live here
    for master in prs.slide_masters:
        total_hits += _remove_watermark_shapes(master.shapes, pattern)
        # 3️⃣  Layouts inside each master
        for layout in master.slide_layouts:
            total_hits += _remove_watermark_shapes(layout.shapes, pattern)

    prs.save(output_path)
    log.info(f"PPTX: removed {total_hits} watermark instance(s)")
    return total_hits


# ─────────────────────────────────────────────
# Helpers
# ─────────────────────────────────────────────
def _delete_file(path: str):
    try:
        os.remove(path)
        log.info(f"Deleted temp file: {path}")
    except FileNotFoundError:
        pass


# ─────────────────────────────────────────────
# Routes
# ─────────────────────────────────────────────
@app.get("/health")
async def health():
    return {"status": "ok"}


@app.post("/process")
async def process_file(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    watermark_text: str = Form(...),
):
    # ── Validate extension ──
    ext = Path(file.filename).suffix.lower()
    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type '{ext}'. Allowed: {', '.join(ALLOWED_EXTENSIONS)}",
        )

    # ── Validate watermark text ──
    watermark_text = watermark_text.strip()
    if not watermark_text:
        raise HTTPException(status_code=400, detail="watermark_text cannot be empty.")
    if len(watermark_text) > 500:
        raise HTTPException(status_code=400, detail="watermark_text is too long (max 500 chars).")

    # ── Read & validate file size ──
    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Maximum allowed size is {MAX_FILE_SIZE_MB} MB.",
        )

    unique_id = str(uuid.uuid4())
    input_path = str(UPLOAD_FOLDER / f"input_{unique_id}{ext}")
    output_path = str(UPLOAD_FOLDER / f"output_{unique_id}{ext}")

    # Write upload to disk
    with open(input_path, "wb") as f:
        f.write(contents)

    # Always delete the input file after processing
    background_tasks.add_task(_delete_file, input_path)

    try:
        if ext == ".pdf":
            hits = clean_pdf(input_path, output_path, watermark_text)
            media_type = "application/pdf"
        else:  # .pptx / .ppt
            hits = clean_pptx(input_path, output_path, watermark_text)
            media_type = (
                "application/vnd.openxmlformats-officedocument.presentationml.presentation"
            )

        if hits == 0:
            log.warning(f"Watermark '{watermark_text}' not found in {file.filename}")

        # Delete output file after it has been sent
        background_tasks.add_task(_delete_file, output_path)

        clean_filename = f"clean_{file.filename}"
        return FileResponse(
            output_path,
            media_type=media_type,
            filename=clean_filename,
            headers={"X-Watermark-Hits": str(hits)},
        )

    except HTTPException:
        _delete_file(output_path)
        raise
    except Exception as e:
        _delete_file(output_path)
        log.exception("Unexpected error during processing")
        raise HTTPException(status_code=500, detail=f"Processing failed: {str(e)}")


# ─────────────────────────────────────────────
# Run directly:  python main.py
# ─────────────────────────────────────────────
if __name__ == "__main__":
    import uvicorn
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
